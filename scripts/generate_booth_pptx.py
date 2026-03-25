#!/usr/bin/env python3
"""
Sky Sentinel — Booth PowerPoint Generator
==========================================
Generates a 10-slide .pptx seamlessly integrating with the official pitch template branding.
Uses the template as the base, keeping its background perfectly visible,
and positions text/UI elements to complement the layout without heavy blocking overlays.
"""

from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
TEMPLATE = ASSETS / "Pitch Presentation -Demo Slide with final template.pptx"
OUT_FILE = ASSETS / "Sky_Sentinel_Booth.pptx"

# App screenshots
SCREENSHOT_DASHBOARD = ASSETS / "screenshot-dashboard.png"
SCREENSHOT_AI_NARRATIVE = ASSETS / "screenshot-ai-narrative.png"
SCREENSHOT_CLUSTERS = ASSETS / "screenshot-clusters.png"
SCREENSHOT_INVESTIGATION = ASSETS / "screenshot-investigation.png"

# Logo
LOGO = ASSETS / "logo-readme.png"

# ─── Template Color Palette & Branding ────────────────────────────────────────
DARK_BLUE = RGBColor(0x09, 0x44, 0x8C)       # #09448C
MED_BLUE = RGBColor(0x1B, 0x85, 0xD2)        # #1B85D2
BRIGHT_BLUE = RGBColor(0x00, 0xB0, 0xF0)     # #00B0F0
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_NAVY = RGBColor(0x05, 0x1A, 0x3A)       # Used for subtle text backing when necessary
NEAR_WHITE = RGBColor(0xF0, 0xF4, 0xF8)

# Dimensions (16:9 Widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Arial"


# ─── XML Shadow Macro ────────────────────────────────────────────────────────
def apply_shadow(shape, opacity_pct=60, blur_pt=4):
    """Add a soft drop shadow to ensure text/shapes pop against complex backgrounds 
    without needing heavy opaque background boxes."""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        
    blur_emu = str(blur_pt * 12700)
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'), {
        'blurRad': blur_emu, 'dist': '38100', 'dir': '5400000', 'algn': 'tl', 'rotWithShape': '0'
    })
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'), {'val': '000000'})
    etree.SubElement(srgbClr, qn('a:alpha'), {'val': str(opacity_pct * 1000)})


# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_text_box(slide, left, top, width, height, text, font_size=24,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, italic=False, shadow=False):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.italic = italic
    p.alignment = alignment
    if shadow:
        apply_shadow(txbox, opacity_pct=75, blur_pt=5)
    return txbox


def add_multiline_text(slide, left, top, width, height, lines, shadow=False):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.get("text", "")
        p.font.size = Pt(line.get("size", 24))
        p.font.bold = line.get("bold", False)
        p.font.color.rgb = line.get("color", WHITE)
        p.font.name = FONT_NAME
        p.font.italic = line.get("italic", False)
        p.alignment = line.get("alignment", PP_ALIGN.LEFT)
        p.space_after = Pt(line.get("spacing_after", 6))
        
    if shadow:
        apply_shadow(txbox, opacity_pct=75, blur_pt=5)
    return txbox


def add_glass_card(slide, left, top, width, height, alpha=60):
    """A subtle, semi-transparent backing so text is readable without hiding the template."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_NAVY
    shape.line.fill.background()
    
    spPr = shape._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    srgbClr = solidFill.find(qn('a:srgbClr'))
    alpha_el = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
    srgbClr.append(alpha_el)
    
    apply_shadow(shape, opacity_pct=40)
    return shape


def build_screenshot(slide, img_path, left, top, width, caption=None):
    """Add a clean screenshot with a subtle border and shadow, letting the template background breathe."""
    border_padding = Inches(0.06)
    img_height = width * 0.58 
    
    # White border frame for contrast against the dark purple/blue background
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - border_padding, top - border_padding, width + border_padding*2, img_height + border_padding*2)
    border.fill.solid()
    border.fill.fore_color.rgb = WHITE
    border.line.fill.background()
    apply_shadow(border, opacity_pct=50, blur_pt=8)

    slide.shapes.add_picture(str(img_path), left, top, width, img_height)

    if caption:
        add_text_box(slide, left, top + img_height + Inches(0.15), width, Inches(0.4),
                     caption, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, italic=True)


def build_header(slide, title_text, color=WHITE):
    """Standardized header placement matching the template's layout flow.
    The template typically has content starting around Y=0.5 to 1.5.
    """
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(11), Inches(0.8), 
                 title_text, font_size=42, bold=True, color=color)


# ─── Slide Builders ──────────────────────────────────────────────────────────

def build_slide_1_title(prs):
    # Layout 3 is perfectly clean for a title
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.4), Inches(2.2), Inches(4.5), Inches(1.8))

    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.0),
                 "AI-Augmented DME Fraud Detection", font_size=44, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.6),
                 "Human-in-the-Loop Pattern Modeling for Medicare Program Integrity",
                 font_size=20, color=BRIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_multiline_text(slide, Inches(1), Inches(6.5), Inches(11.333), Inches(0.8), [
        {"text": "ACT-IAC AI Hackathon 2026: AI in Action · Sky Solutions", "size": 16, "bold": True, "alignment": PP_ALIGN.CENTER}
    ])


def build_slide_2_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "The Problem")

    add_text_box(slide, Inches(0.5), Inches(2.0), Inches(11.333), Inches(1.8),
                 "$60 Billion+", font_size=100, bold=True, color=BRIGHT_BLUE)

    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(10.333), Inches(0.7),
                 "Lost to Medicare Fraud, Waste & Abuse Every Year", font_size=32, bold=True, color=DARK_BLUE)

    add_multiline_text(slide, Inches(0.5), Inches(4.8), Inches(11), Inches(1.5), [
        {"text": "Modern fraud networks fragment activity across dozens of shell companies.", "size": 24, "color": DARK_BLUE},
        {"text": "Legacy systems catch isolated actors.", "size": 24, "color": DARK_BLUE},
        {"text": "Sky Sentinel catches the entire coordinated network.", "size": 28, "bold": True, "color": BRIGHT_BLUE, "spacing_before": 20}
    ])


def build_slide_3_ensemble(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "Three-Layer AI Detection Pipeline")

    cx = Inches(0.5)
    cw = Inches(4.0)
    gap = Inches(0.2)
    cy = Inches(2.2)

    cards = [
        {"icon": "🧠", "t": "Ensemble ML Detection", "c": [
            "Isolation Forest: Outliers", "Z-Score: Peer deviations", "Time-Series: Ramp-ups", "Parallel execution"]},
        {"icon": "🔍", "t": "Behavioral Clustering", "c": [
            "DBSCAN behavioral grouping", "Discovers networks automatically", "No pre-defined rules needed", "Caught 6 rings in Gold Rush"]},
        {"icon": "💬", "t": "Two-Stage LLM Pipeline", "c": [
            "Stage 1: Encoder proxy", "Stage 2: Decoder narrative", "Plain-English explanations", "Vendor-agnostic architecture"]}
    ]

    for i, card in enumerate(cards):
        x = cx + (cw + gap)*i
        add_glass_card(slide, x, cy, cw, Inches(4.5), alpha=60)
        
        add_text_box(slide, x, cy + Inches(0.2), cw, Inches(0.8), card["icon"], font_size=50, alignment=PP_ALIGN.CENTER, shadow=False)
        add_text_box(slide, x, cy + Inches(1.2), cw, Inches(1.0), card["t"], font_size=24, bold=True, color=BRIGHT_BLUE, alignment=PP_ALIGN.CENTER, shadow=False)
        
        lines = [{"text": f"• {item}", "size": 18, "spacing_after": 12} for item in card["c"]]
        add_multiline_text(slide, x + Inches(0.3), cy + Inches(2.0), cw - Inches(0.6), Inches(2.0), lines, shadow=False)


def build_slide_4_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "Command Center Dashboard")
    if SCREENSHOT_DASHBOARD.exists():
        build_screenshot(slide, SCREENSHOT_DASHBOARD, Inches(1.0), Inches(1.6), Inches(11.333), "Live Monitoring · 347 Suppliers · 5,500+ Claims")


def build_slide_5_explainable(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "Explainable AI Narratives")

    if SCREENSHOT_AI_NARRATIVE.exists():
        build_screenshot(slide, SCREENSHOT_AI_NARRATIVE, Inches(0.5), Inches(1.6), Inches(7.5))

    rx = Inches(8.3)
    rw = Inches(4.5)
    
    add_glass_card(slide, rx, Inches(1.6), rw, Inches(4.3), alpha=60)
    add_multiline_text(slide, rx + Inches(0.3), Inches(1.9), rw - Inches(0.6), Inches(4.0), [
        {"text": "Every Alert Includes:", "size": 24, "bold": True, "color": BRIGHT_BLUE, "spacing_after": 16},
        {"text": "✦ Plain-English narrative", "size": 18, "spacing_after": 12},
        {"text": "✦ HCPCS codes & dollar amounts", "size": 18, "spacing_after": 12},
        {"text": "✦ 6-factor decomposable score", "size": 18, "spacing_after": 12},
        {"text": "✦ Key concerns & evidence", "size": 18, "spacing_after": 30},
        {"text": "No black boxes. Humans decide.", "size": 20, "bold": True, "color": WHITE, "alignment": PP_ALIGN.CENTER}
    ], shadow=False)


def build_slide_6_clusters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "Discovering Coordinated Fraud Rings")
    if SCREENSHOT_CLUSTERS.exists():
        build_screenshot(slide, SCREENSHOT_CLUSTERS, Inches(1.0), Inches(1.6), Inches(11.333), "DBSCAN Clustering · 6 Rings Detected · Patterns Invisible to Single-Supplier Checking")


def build_slide_7_human_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "Human-In-The-Loop Controls")
    if SCREENSHOT_INVESTIGATION.exists():
        build_screenshot(slide, SCREENSHOT_INVESTIGATION, Inches(1.0), Inches(1.6), Inches(11.333), "6-Dimension Tunable Weights · Hypothesis Testing · Saved Configurations")


def build_slide_8_gold_rush(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    build_header(slide, "Real-World Validation: Operation Gold Rush")

    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5), "The largest healthcare fraud case by dollar amount ever charged by the DOJ.", font_size=20, italic=True)

    def draw_stat(x, y, num, text):
        add_glass_card(slide, x, y, Inches(3.8), Inches(1.8), alpha=50)
        add_text_box(slide, x, y+Inches(0.1), Inches(3.8), Inches(1.0), num, font_size=54, bold=True, color=BRIGHT_BLUE, alignment=PP_ALIGN.CENTER, shadow=False)
        add_text_box(slide, x, y+Inches(1.0), Inches(3.8), Inches(0.8), text, font_size=18, alignment=PP_ALIGN.CENTER, shadow=False)

    draw_stat(Inches(0.5),  Inches(2.5), "$10.6B", "Fraudulent Claims")
    draw_stat(Inches(4.7), Inches(2.5), "1M+", "Stolen Identities")
    draw_stat(Inches(8.9), Inches(2.5), "50", "States Affected")
    
    draw_stat(Inches(0.5),  Inches(4.8), "6", "Fraud Rings Detected")
    draw_stat(Inches(4.7), Inches(4.8), "34", "Linked Shell Suppliers")
    draw_stat(Inches(8.9), Inches(4.8), "100%", "Networks Discovered AI")


def build_slide_9_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    # For a brighter 2-column layout (Blue), headers might want to be darker if background is white/light
    build_header(slide, "Built for Production")

    c1 = Inches(0.5)
    c2 = Inches(6.9)
    cw = Inches(5.9)
    ch = Inches(5.0)

    add_glass_card(slide, c1, Inches(1.8), cw, ch, alpha=50)
    add_multiline_text(slide, c1+Inches(0.3), Inches(2.2), cw-Inches(0.6), ch, [
        {"text": "Data & AI Architecture", "size": 26, "bold": True, "color": BRIGHT_BLUE, "spacing_after": 20},
        {"text": "▸ Live CMS Medicare Supplier API", "size": 20, "spacing_after": 14},
        {"text": "▸ 347 suppliers + 5,500+ claims", "size": 20, "spacing_after": 14},
        {"text": "▸ Isolation Forest + DBSCAN", "size": 20, "spacing_after": 14},
        {"text": "▸ Vendor-agnostic LLM tiering", "size": 20, "spacing_after": 14},
        {"text": "▸ Decomposable risk scoring", "size": 20}
    ], shadow=False)

    add_glass_card(slide, c2, Inches(1.8), cw, ch, alpha=50)
    add_multiline_text(slide, c2+Inches(0.3), Inches(2.2), cw-Inches(0.6), ch, [
        {"text": "Deployment Ready", "size": 26, "bold": True, "color": BRIGHT_BLUE, "spacing_after": 20},
        {"text": "▸ Docker single-container deploy", "size": 20, "spacing_after": 14},
        {"text": "▸ React 19 + FastAPI + SQLite", "size": 20, "spacing_after": 14},
        {"text": "▸ JWT + RBAC (3 distinct roles)", "size": 20, "spacing_after": 14},
        {"text": "▸ 25+ Complete REST API endpoints", "size": 20, "spacing_after": 14},
        {"text": "▸ Scalable to PostgreSQL", "size": 20}
    ], shadow=False)


def build_slide_10_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.4), Inches(2.0), Inches(4.5), Inches(1.8))

    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.333), Inches(1.0),
                 "Ask Us for a Live Demo!", font_size=54, bold=True, color=BRIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.6),
                 "See how AI amplifies human expertise to catch coordinated fraud", font_size=24, alignment=PP_ALIGN.CENTER)

    add_multiline_text(slide, Inches(1), Inches(6.5), Inches(11.333), Inches(0.8), [
        {"text": "Team 1 · Sky Solutions", "size": 18, "bold": True, "alignment": PP_ALIGN.CENTER, "spacing_after": 4},
        {"text": "ACT-IAC AI Hackathon 2026", "size": 14, "alignment": PP_ALIGN.CENTER}
    ])


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    print("🚀 Sky Sentinel Booth PowerPoint Generator (Seamless Template Integration)")
    prs = Presentation(str(TEMPLATE))

    # Strip existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    builders = [
        ("Title", build_slide_1_title), ("The Problem", build_slide_2_problem),
        ("Ensemble AI", build_slide_3_ensemble), ("Command Center", build_slide_4_dashboard),
        ("Explainable AI", build_slide_5_explainable), ("Fraud Networks", build_slide_6_clusters),
        ("Human-in-the-Loop", build_slide_7_human_loop), ("Operation Gold Rush", build_slide_8_gold_rush),
        ("Tech & Deployment", build_slide_9_tech), ("Call to Action", build_slide_10_cta)
    ]

    for i, (n, b) in enumerate(builders, 1):
        print(f"[{i}/10] {n}")
        b(prs)

    prs.save(str(OUT_FILE))
    print(f"✅ Saved: {OUT_FILE} ({OUT_FILE.stat().st_size/1024/1024:.1f}MB)")

if __name__ == "__main__":
    main()
